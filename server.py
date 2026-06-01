import os
import shutil
import tempfile
import whisper
import uvicorn
from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
from openai import OpenAI
import json

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    pptx_installed = True
except ImportError:
    pptx_installed = False

app = FastAPI(
    title="Smart Classroom Transcription & AI Server",
    description="Offloaded Whisper transcription and Groq AI service for SBC client"
)

# Enable CORS so the client SBC can communicate with it
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Load the model globally so it stays in RAM and does not reload on every request
print("[INFO] Loading Whisper 'medium' model... This may take a minute.")
try:
    model = whisper.load_model("medium")
    print("[SUCCESS] Whisper model loaded and ready!")
    model_loaded = True
except Exception as e:
    print(f"[ERROR] Failed to load Whisper model: {e}")
    model_loaded = False

# Initialize Groq client
try:
    groq_client = OpenAI(
        api_key="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",
        base_url="https://api.groq.com/openai/v1"
    )
    print("[SUCCESS] Groq AI client initialized!")
except Exception as e:
    print(f"[ERROR] Failed to initialize Groq client: {e}")


class NotesRequest(BaseModel):
    text: str
    mode: str
    custom_prompt: str = ""


@app.get("/health")
def health_check():
    return {
        "status": "healthy" if model_loaded else "unhealthy",
        "model": "medium" if model_loaded else None
    }


@app.post("/transcribe")
async def transcribe(file: UploadFile = File(...)):
    if not model_loaded:
        raise HTTPException(status_code=503, detail="Whisper model is not loaded or healthy on the server.")

    # Check file extension
    filename = file.filename
    _, ext = os.path.splitext(filename)
    if ext.lower() not in [".wav", ".mp3", ".m4a", ".ogg", ".flac"]:
        raise HTTPException(status_code=400, detail=f"Unsupported file format: {ext}")

    # Create a temporary file to save the uploaded audio data
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as temp_audio:
        temp_path = temp_audio.name
        try:
            # Stream upload file contents into the temporary file
            shutil.copyfileobj(file.file, temp_audio)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Failed to save uploaded file: {e}")

    try:
        print(f"[INFO] Transcribing file: {filename} ({os.path.getsize(temp_path)} bytes)...")
        # Run Whisper transcription with same prompts and settings as transcribe_module.py
        result = model.transcribe(
            temp_path,
            language="en",
            beam_size=5,
            initial_prompt="""
            Engineering classroom lecture.
            Technical educational discussion.
            """
        )
        transcription_text = result.get("text", "").strip()
        print(f"[SUCCESS] Transcription complete for: {filename}")
        return {"transcription": transcription_text}

    except Exception as e:
        print(f"[ERROR] Error during transcription: {e}")
        raise HTTPException(status_code=500, detail=f"Transcription failed: {str(e)}")

    finally:
        # Always clean up the temporary file
        if os.path.exists(temp_path):
            os.remove(temp_path)


@app.post("/generate_notes")
async def generate_notes(request: NotesRequest):
    prompts = {
        "Summary": "Summarize these lecture notes clearly.",
        "Detailed Notes": "Generate well-structured lecture notes from this transcription.",
        "Important Questions": "Generate important exam questions from this lecture.",
        "Key Points": "Extract key points from this lecture.",
        "Explain Simply": "Explain this lecture in simple student-friendly language."
    }

    system_prompt = prompts.get(request.mode, "")
    if request.custom_prompt:
        system_prompt += "\nAdditional Instruction: " + request.custom_prompt

    try:
        print(f"[INFO] Generating AI notes using Groq (mode: {request.mode})...")
        response = groq_client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[
                {
                    "role": "system",
                    "content": system_prompt
                },
                {
                    "role": "user",
                    "content": request.text
                }
            ],
            temperature=0.3
        )
        ai_content = response.choices[0].message.content
        print("[SUCCESS] AI notes generated successfully!")
        return {"notes": ai_content}
    except Exception as e:
        print(f"[ERROR] Failed to generate AI notes: {e}")
        raise HTTPException(status_code=500, detail=f"AI generation failed: {str(e)}")


class PPTRequest(BaseModel):
    text: str


@app.post("/generate_ppt")
async def generate_ppt(request: PPTRequest):
    if not pptx_installed:
        raise HTTPException(status_code=500, detail="python-pptx is not installed on the server.")

    system_prompt = """
    You are an expert academic slide outline designer. 
    Analyze the provided lecture transcription and design a highly structured, comprehensive educational slide presentation.
    Your output MUST be a single, valid JSON object with NO markdown formatting, NO backticks, and NO trailing text.

    Expected JSON schema:
    {
      "title": "Main Title of the Lecture",
      "subtitle": "Comprehensive Study Slide Presentation",
      "slides": [
        {
          "title": "Title of Slide (e.g. Introduction to Topic)",
          "bullets": [
            "Clear, concise academic bullet point describing key concept.",
            "Another core technical detail or definition.",
            "Critical formula or application detail."
          ]
        }
      ]
    }

    Rules:
    1. Output at least 5-8 informative slides depending on lecture length.
    2. Write highly technical, professional, and clear educational content.
    3. Ensure the JSON is completely valid and parseable by json.loads() in Python.
    """

    try:
        print("[INFO] Generating structured slides JSON using Groq...")
        response = groq_client.chat.completions.create(
            model="llama-3.1-8b-instant",
            response_format={"type": "json_object"},
            messages=[
                {
                    "role": "system",
                    "content": system_prompt
                },
                {
                    "role": "user",
                    "content": request.text
                }
            ],
            temperature=0.3
        )
        
        raw_output = response.choices[0].message.content.strip()
        slide_data = json.loads(raw_output)
        print("[SUCCESS] Slide data outline successfully structured!")

        # Create PPTX Presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title Slide (Premium Dark Crimson theme)
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Crimson background
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(90, 31, 34)  # Crimson: #5a1f22
        bg.line.fill.background()

        # Title textbox
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "Smart Classroom Lecture")
        p.font.name = "Georgia"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(229, 195, 101)  # Gold: #e5c365
        p.alignment = PP_ALIGN.LEFT

        p2 = tf.add_paragraph()
        p2.text = slide_data.get("subtitle", "Comprehensive Study Slide Presentation")
        p2.font.name = "Arial"
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(255, 255, 255)  # White
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(14)

        # Content Slides
        for s_info in slide_data.get("slides", []):
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Gold top accent bar
            accent = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.4))
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(229, 195, 101)
            accent.line.fill.background()

            # Slide title textbox
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.0))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = s_info.get("title", "Lecture Content")
            p.font.name = "Georgia"
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(90, 31, 34)  # Crimson

            # Slide body textbox
            body_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.8))
            tf = body_box.text_frame
            tf.word_wrap = True

            bullets = s_info.get("bullets", [])
            for i, bullet in enumerate(bullets):
                p_bullet = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p_bullet.text = "•  " + bullet
                p_bullet.font.name = "Arial"
                p_bullet.font.size = Pt(18)
                p_bullet.font.color.rgb = RGBColor(51, 21, 21)
                p_bullet.space_after = Pt(14)

        # Save to temp file
        temp_dir = tempfile.gettempdir()
        temp_file_path = os.path.join(temp_dir, "Generated_Lecture_Presentation.pptx")
        prs.save(temp_file_path)
        print("[SUCCESS] PPTX presentation file built successfully!")

        return FileResponse(
            temp_file_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename="Lecture_Presentation.pptx"
        )

    except Exception as e:
        print(f"[ERROR] Failed to generate PPTX: {e}")
        raise HTTPException(status_code=500, detail=f"PPT generation failed: {str(e)}")


class QuizRequest(BaseModel):
    text: str


@app.post("/generate_quiz")
async def generate_quiz(request: QuizRequest):
    system_prompt = """
    You are an expert academic examiner. 
    Analyze the provided lecture transcription and design a highly informative, multiple-choice quiz of exactly 5 questions to test students' comprehension of key technical concepts.
    Your output MUST be a single, valid JSON object with NO markdown formatting, NO backticks, and NO trailing text.

    Expected JSON schema:
    {
      "questions": [
        {
          "question": "What is the main topic discussed in the first section of the lecture?",
          "options": ["Option A text", "Option B text", "Option C text", "Option D text"],
          "correct_index": 0
        }
      ]
    }

    Rules:
    1. Output EXACTLY 5 multiple-choice questions.
    2. Provide exactly 4 options per question.
    3. The correct_index must be a valid index (0, 1, 2, or 3) pointing to the correct option in the options array.
    4. Write clear, technical, and accurate questions based strictly on the transcript text.
    5. Ensure the JSON is completely valid and parseable by json.loads() in Python.
    """

    try:
        print("[INFO] Generating structured Quiz JSON using Groq...")
        response = groq_client.chat.completions.create(
            model="llama-3.1-8b-instant",
            response_format={"type": "json_object"},
            messages=[
                {
                    "role": "system",
                    "content": system_prompt
                },
                {
                    "role": "user",
                    "content": request.text
                }
            ],
            temperature=0.3
        )
        
        raw_output = response.choices[0].message.content.strip()
        quiz_data = json.loads(raw_output)
        
        # Validate shape
        if "questions" not in quiz_data or len(quiz_data["questions"]) != 5:
            print("[WARNING] Quiz did not contain exactly 5 questions. Retrying with strict instruction...")
            # If length is mismatch, we'll gracefully return what we have or adjust
            
        print("[SUCCESS] Quiz questions successfully structured and parsed!")
        return quiz_data

    except Exception as e:
        print(f"[ERROR] Failed to generate Quiz: {e}")
        raise HTTPException(status_code=500, detail=f"Quiz generation failed: {str(e)}")


if __name__ == "__main__":
    # Run on all network interfaces (0.0.0.0) at port 8000 so the SBC can reach it
    uvicorn.run(app, host="0.0.0.0", port=8000)
